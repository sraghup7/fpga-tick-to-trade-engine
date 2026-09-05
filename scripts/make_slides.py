#!/usr/bin/env python3
"""Generate FPGA_Tick_to_Trade.pptx from project_explained_simple.md content."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------- constants
SW, SH = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT = RGBColor(0xE8, 0x8C, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(box, color)
    box.shadow.inherit = False
    return box


def add_text(slide, left, top, width, height, text, size=16, color=DARK,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=15, gap=6,
                color=DARK, bullet_color=None):
    """items: list of (level, text). level 0 bullet, 1 sub, -1 plain."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ""
        p.space_after = Pt(gap)
        p.level = max(level, 0)
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size - (2 if level else 0))
        run.font.color.rgb = color
        if level < 0:
            run.font.bold = True
            run.font.size = Pt(size + 2)
        elif level == 0 and bullet_color is not None:
            run.font.color.rgb = bullet_color
    return box


def add_table(slide, left, top, width, rows, cols, data, col_widths=None,
              font_size=12, row_height=0.34, header_size=12, header_fill=NAVY):
    shape = slide.shapes.add_table(rows, cols, left, top, width,
                                   Inches(row_height * rows))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.text = ""
                run = p.add_run()
                run.text = str(val)
                run.font.name = FONT
                if r == 0:
                    run.font.size = Pt(header_size)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
    return tbl


def content_slide(step_label, title):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    if step_label:
        add_text(slide, Inches(0.5), Inches(0.18), Inches(2.0), Inches(0.7),
                 step_label, size=16, color=ACCENT, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(2.3), Inches(0.12), Inches(10.5), Inches(0.8),
                 title, size=26, color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    else:
        add_text(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.8),
                 title, size=26, color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, Inches(1.0), SW, Inches(0.06), ACCENT)
    return slide


def footer(slide, num):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(8.0), Inches(0.3),
             "FPGA Tick-to-Trade Simulator", size=9, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.12), Inches(0.8), Inches(0.3),
             str(num), size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- slide 1: title
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SW, SH, NAVY)
add_rect(slide, 0, Inches(4.9), SW, Inches(0.06), ACCENT)
add_text(slide, Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.6),
         "FPGA Tick-to-Trade", size=54, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(3.3), Inches(10.9), Inches(1.2),
         "A simulated high-frequency trading engine on silicon — "
         "from price message to order in ~176 nanoseconds, with perfectly "
         "fixed latency", size=18, color=RGBColor(0xD6, 0xE4, 0xF0),
         align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.5),
         "9 steps, explained simply  |  Spec-only, pre-implementation",
         size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(6.2), Inches(10.9), Inches(0.9),
         "Source: fpga_tick_to_trade_master_spec.md v2.0 + ml_engineer_brief.md  ·  "
         "Target: Artix-7 XC7A35T-2FGG484I @ 125 MHz  ·  Vivado/Vitis HLS 2023.x",
         size=12, color=RGBColor(0x9C, 0xB4, 0xC8), align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- slide 2: overview
slide = content_slide(None, "The Project in One Glance")
add_bullets(slide, Inches(0.7), Inches(1.5), Inches(6.0), Inches(5.4), [
    (0, "A super-fast simulated trading system on an FPGA chip"),
    (0, "No real money, no real exchange — a portfolio demo"),
    (0, "Goal: prove you can make trading decisions in ~176 ns "
        "with perfectly fixed latency"),
    (0, "Verified with 1,000,000+ messages, never dropping or losing one"),
    (0, "Every trade checked by a tiny AI before it leaves"),
], size=16, gap=10)
steps = [
    "1. The Problem — why speed matters",
    "2. The Mail Format — how price messages look",
    "3. The Ears — how the FPGA receives mail",
    "4. The Notebook — tracking top of book",
    "5. The Clues — extracting 8 market hints",
    "6. The AI Brain — predicting a risky trade",
    "7. The Trading Rule — deciding to buy/sell",
    "8. The 9 Security Guards — blocking danger",
    "9. The Exit & Stopwatch — send + prove speed",
]
box = slide.shapes.add_textbox(Inches(7.1), Inches(1.5), Inches(5.5), Inches(5.4))
tf = box.text_frame
tf.word_wrap = True
for i, s in enumerate(steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(7)
    run = p.add_run()
    run.text = s
    run.font.name = FONT
    run.font.size = Pt(14)
    run.font.color.rgb = DARK
footer(slide, 2)

# ---------------------------------------------------------------- slide 3: step 1
slide = content_slide("STEP 1", "The Problem — Why Speed Matters")
add_bullets(slide, Inches(0.7), Inches(1.35), Inches(12.0), Inches(2.6), [
    (0, "The stock market is a live auction — everyone hears the price at the "
        "same second, and the first to shout \"I'LL BUY!\" gets the deal. "
        "One microsecond slow = it's gone."),
    (0, "A normal computer reacts unpredictably: sometimes 2 µs, sometimes 50 µs "
        "(it gets slow exactly when the market is most busy)."),
    (0, "An FPGA is a fixed assembly line, not a multitasking person — every "
        "price message takes the exact same physical path, so time is always "
        "identical, e.g. exactly 22 ticks ≈ 176 nanoseconds."),
], size=16, gap=10)
add_rect(slide, Inches(0.7), Inches(4.0), Inches(12.0), Inches(0.9), LIGHT)
add_text(slide, Inches(1.0), Inches(4.22), Inches(11.4), Inches(0.6),
         "Project goal: build the assembly line that listens, decides, "
         "double-checks with AI, and sends a fake order back — all in a fixed, "
         "tiny amount of time.",
         size=15, color=NAVY, bold=True)
add_bullets(slide, Inches(0.7), Inches(5.1), Inches(12.0), Inches(1.8), [
    (0, "Listen to fake price messages from a laptop"),
    (0, "Decide if it's a good time to trade"),
    (0, "Double-check with an AI if it's risky"),
    (0, "Send a fake order back — and prove no message was ever lost, "
        "even over 1,000,000 messages in a row"),
], size=15, gap=5)
footer(slide, 3)

# ---------------------------------------------------------------- slide 4: step 2
slide = content_slide("STEP 2", "The Mail Format — The 16-Byte Postcard")
add_bullets(slide, Inches(0.7), Inches(1.25), Inches(6.0), Inches(2.2), [
    (0, "Every price update is a standard postcard: exactly 16 bytes, 6 boxes, "
        "always in the same place (big-endian / network byte order)."),
    (0, "Because the size never changes, the FPGA just counts — it never has "
        "to guess where a message ends."),
], size=14, gap=8)
add_table(slide, Inches(0.7), Inches(2.6), Inches(7.4), 7, 3,
          [["Box", "Means", "Example"],
           ["Type (1B)", "Kind of news", "0x01 update · 0x02 trade · 0x03 clear · 0xFF heartbeat"],
           ["Symbol (1B)", "Which item", "1 = Apples, 2 = Oranges (watch 4)"],
           ["Side (1B)", "Which side", "0 = BUY, 1 = SELL (trade: aggressor)"],
           ["Price (4B)", "Price in cents", "10001 = $100.01 (integers only)"],
           ["Qty (4B)", "How many", "50 apples · 0 = side empty"],
           ["Seq # (4B)", "Counting 1,2,3…", "1,2,4 ⇒ #3 was lost"]],
          col_widths=[1.5, 1.7, 4.2], font_size=11, row_height=0.62, header_size=12)
add_bullets(slide, Inches(8.5), Inches(1.25), Inches(4.3), Inches(5.6), [
    (1, "Two tricks"),
    (0, "Packing: one Ethernet frame holds 1–88 postcards (88×16 = 1408 bytes), "
        "streamed back-to-back"),
    (0, "Two envelopes: Mode 0 raw Ethernet (0x88B5) = training wheels; "
        "Mode 1 real UDP/IPv4 (port 60000) = the real road"),
    (0, "Reply postcard is also 16 bytes: Type 0x10/0x11/0x12, symbol, side, "
        "reject reason, price, qty, trigger_seq, latency_cyc"),
    (0, "LINK_MODE selects the envelope — the assembly line doesn't care"),
], size=13, gap=7)
footer(slide, 4)

# ---------------------------------------------------------------- slide 5: step 3
slide = content_slide("STEP 3", "The Ears — How the FPGA Receives Mail")
rows = [
    ["Station", "Module", "Job"],
    ["A", "rgmii_to_gmii", "Translator: 4-lane dual-sided highway (RGMII DDR @125 MHz) → orderly 1-byte-per-tick belt"],
    ["B", "eth_mac_rx", "Mailroom: strips preamble/SFD, verifies FCS; damaged envelopes trashed + counted (err_fcs)"],
    ["C", "frame_classifier", "Guard: Mode 0 → EtherType 0x88B5; Mode 1 → IPv4 + IHL=5 + checksum + UDP port 60000; per-reason counters"],
    ["D", "md_parser", "Opener: cuts payload into exact 16-byte postcards; length not divisible by 16 → trash (err_frame_len)"],
]
add_table(slide, Inches(0.7), Inches(1.35), Inches(12.0), 5, 3, rows,
          col_widths=[1.0, 2.6, 8.4], font_size=12, row_height=0.62, header_size=12)
add_bullets(slide, Inches(0.7), Inches(5.15), Inches(12.0), Inches(1.8), [
    (0, "At the very last byte entering Station D, we stamp the time — that is "
        "the \"start\" of the latency stopwatch, carried with the message."),
    (0, "Everything up to here is about getting mail cleanly without dropping "
        "anything, at minimum inter-frame gap for 10,000 frames."),
], size=14, gap=8)
footer(slide, 5)

# ---------------------------------------------------------------- slide 6: step 4
slide = content_slide("STEP 4", "The Notebook — Tracking Top of Book")
add_bullets(slide, Inches(0.7), Inches(1.3), Inches(6.1), Inches(5.5), [
    (0, "A tiny notebook with 4 pages (NUM_SYMBOLS=4): Apples, Oranges, "
        "Bananas, Grapes"),
    (0, "Each page holds only the top of book: best BUY (bid) price+qty and "
        "best SELL (ask) price+qty, each with a valid bit"),
    (0, "Spread = ask − bid · Mid = (bid+ask) >> 1"),
    (0, "Book starts invalid on reset; crossed (bid ≥ ask, both valid) flags "
        "cnt_crossed and later guards block trading"),
], size=14, gap=8)
rows = [
    ["Station", "Module", "Job"],
    ["E", "symbol_filter", "Bouncer: 4-entry CAM picks only our symbols; else discard + cnt_filtered"],
    ["F", "seq_monitor", "Missed-mail detector: gap → sticky seq_gap + stale (block until resync); duplicate → drop; 10 ms staleness timer"],
    ["G", "tob_engine", "Notebook keeper: 1-cycle replace per side; qty=0 clears valid; 0x03 clears book; 0x02/0xFF don't touch book"],
]
add_table(slide, Inches(7.0), Inches(1.3), Inches(5.8), 4, 3, rows,
          col_widths=[1.0, 2.3, 2.5], font_size=10.5, row_height=0.85, header_size=11)
add_text(slide, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.5),
         "All later steps read from this notebook.", size=14, color=NAVY, bold=True)
footer(slide, 6)

# ---------------------------------------------------------------- slide 7: step 5
slide = content_slide("STEP 5", "The Clues — 8 Market Hints (Features)")
rows = [
    ["Clue", "Name", "Formula", "Plain English"],
    ["F0", "Spread", "a_t − b_t", "Is the buyer/seller gap wide or narrow?"],
    ["F1", "Mid delta", "m_t − m_(t−1)", "Did the middle price move up/down?"],
    ["F2", "Book imbalance", "q_b − q_a", "More buyers or more sellers?"],
    ["F3", "Bid-size change", "q_b,t − q_b,t−1", "Did buyers grow/shrink?"],
    ["F4", "Ask-size change", "q_a,t − q_a,t−1", "Did sellers grow/shrink?"],
    ["F5", "Update rate", "updates in last W events", "How busy is the market? (W=16)"],
    ["F6", "Last trade side", "buy +1 / sell −1 / none 0", "Who started the last trade?"],
    ["F7", "Volatility proxy", "Σ|mid deltas| (last W)", "How much did price jitter?"],
]
add_table(slide, Inches(0.7), Inches(1.3), Inches(8.6), 9, 4, rows,
          col_widths=[0.7, 1.9, 3.2, 2.8], font_size=11, row_height=0.5, header_size=12)
add_bullets(slide, Inches(9.6), Inches(1.3), Inches(3.2), Inches(5.5), [
    (0, "Computed after every book update in 1 tick — add/subtract/shift/compare only"),
    (0, "Wide spread = low interest; narrow = crowded, active market"),
    (0, "Squeezer (feature_normalizer) maps raw clues to int8 −128..127 via "
        "offset + arithmetic shift, saturating not wrapping"),
    (0, "Result: 8 tiny numbers — the exact AI input"),
], size=12.5, gap=8)
footer(slide, 7)

# ---------------------------------------------------------------- slide 8: step 6
slide = content_slide("STEP 6", "The AI Brain — Predicting a Risky Trade")
add_bullets(slide, Inches(0.7), Inches(1.3), Inches(6.2), Inches(5.5), [
    (0, "A tiny AI that answers ONE question: \"if we quote now, will the price "
        "move against us and we get adversely selected?\""),
    (0, "Adverse selection = you sell at $10, price crashes to $8, you're stuck"),
    (0, "Score:  z = b + w0·x0 + w1·x1 + … + w7·x7"),
    (0, "8 weights learned in training; int32 exact, no silent truncation"),
    (0, "It does NOT buy or sell — it only vetoes"),
], size=14, gap=9)
add_rect(slide, Inches(7.2), Inches(1.3), Inches(5.4), Inches(2.2), LIGHT)
add_text(slide, Inches(7.6), Inches(1.5), Inches(4.6), Inches(1.9),
         "Hysteresis (two lines, no flicker)\n\n"
         "z ≥ T_high → Risky = 1 (block)\n"
         "z ≤ T_low  → Safe = 0 (allow)\n"
         "in between → hold previous",
         size=13, color=DARK)
add_rect(slide, Inches(7.2), Inches(3.7), Inches(5.4), Inches(1.5), LIGHT)
add_text(slide, Inches(7.6), Inches(3.85), Inches(4.6), Inches(1.3),
         "Fail-safe (FR-31): invalid / crossed / gap / stale data "
         "⇒ adverse_risk forced to 1 — better to block than trade on garbage.",
         size=12.5, color=DARK)
add_bullets(slide, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.4), [
    (0, "Runs in parallel with the trading rule — the rule's order intent is "
        "delayed through a fixed-depth [ALIGN] shift register so both arrive at "
        "the risk engine the same cycle, every cycle (NFR-14)."),
], size=13, gap=8)
footer(slide, 8)

# ---------------------------------------------------------------- slide 9: step 7
slide = content_slide("STEP 7", "The Trading Rule — Deciding to Buy or Sell")
add_bullets(slide, Inches(0.7), Inches(1.3), Inches(5.9), Inches(5.6), [
    (0, "A simple, fixed, deterministic rule — not AI. The AI only vetoes."),
    (0, "Wakes up after every book-modifying update and asks: should we trade?"),
    (0, "Cheap math only: barrel-shift, compare, subtract — no DSP on signal path"),
    (0, "Takes just 1 tick; intent delayed via [ALIGN] to meet the AI verdict "
        "at the same cycle"),
], size=14, gap=9)
add_rect(slide, Inches(6.9), Inches(1.3), Inches(5.7), Inches(2.75), LIGHT)
add_text(slide, Inches(7.2), Inches(1.45), Inches(5.1), Inches(2.4),
         "BUY (FR-35): buy order iff\n"
         "• both sides valid, not crossed\n"
         "• spread ≥ cfg_min_spread (2 ticks)\n"
         "• bid_qty > ask_qty << 1  (way more buyers)\n"
         "→ Buy cfg_order_qty (100) at ask_price",
         size=13, color=DARK)
add_rect(slide, Inches(6.9), Inches(4.25), Inches(5.7), Inches(2.75), LIGHT)
add_text(slide, Inches(7.2), Inches(4.4), Inches(5.1), Inches(2.4),
         "SELL (FR-36): mirror image\n"
         "• both sides valid, not crossed\n"
         "• spread ≥ cfg_min_spread\n"
         "• ask_qty > bid_qty << 1  (way more sellers)\n"
         "→ Sell cfg_order_qty at bid_price",
         size=13, color=DARK)
add_text(slide, Inches(0.7), Inches(6.35), Inches(5.9), Inches(0.9),
         "Both true at once is impossible by construction — if ever detected, "
         "no order + err_signal_conflict (FR-37).",
         size=12, color=GRAY)
footer(slide, 9)

# ---------------------------------------------------------------- slide 10: step 8
slide = content_slide("STEP 8", "The 9 Security Guards — Blocking Dangerous Orders")
rows = [
    ["Gate", "Blocks if…"],
    ["0x01 Kill Switch", "Operator halt latched — overrides everything (CSR clear to resume)"],
    ["0x02 Max Size", "order_qty > 500 (no fat-finger)"],
    ["0x03 Max Position", "|position ± qty| > 1000 exposure bound"],
    ["0x04 Price Band", "|order price − mid| > 50 ticks (corrupt price)"],
    ["0x05 Stale Data", "no update for 10 ms"],
    ["0x06 Sequence Gap", "sticky seq_gap — book unreliable until resync"],
    ["0x07 Crossed/Locked", "bid ≥ ask both valid — physically impossible"],
    ["0x08 Throttle", "token bucket empty (8 tokens, 100 µs refill)"],
    ["0x09 Adverse Selection (ML)", "adverse_risk = 1 → block, or reduce size (cfg_ml_action)"],
]
add_table(slide, Inches(0.7), Inches(1.3), Inches(8.3), 10, 2, rows,
          col_widths=[2.6, 5.7], font_size=11, row_height=0.5, header_size=12)
add_bullets(slide, Inches(9.3), Inches(1.3), Inches(3.4), Inches(5.5), [
    (0, "All 9 check in parallel in a single cycle (FR-42) — any raised hand blocks"),
    (0, "No path bypasses the risk engine (FR-41)"),
    (0, "Every block increments that guard's counter"),
    (0, "If 2 fire at once: both counted, reject_reason = lowest-numbered gate "
        "(hard gates 0x01–0x08 dominate ML 0x09)"),
    (0, "ML 0x09 is advisory — never overrides hard gates"),
    (0, "All 9 pass ⇒ intent becomes a real order"),
], size=13, gap=8)
footer(slide, 10)

# ---------------------------------------------------------------- slide 11: step 9 exit
slide = content_slide("STEP 9", "The Exit — Sending the Order")
add_bullets(slide, Inches(0.7), Inches(1.3), Inches(6.0), Inches(5.5), [
    (0, "order_builder + eth_mac_tx encode the order with correct framing "
        "for the active LINK_MODE (incl. IPv4 checksum in Mode 1)"),
    (0, "~1 tick from risk pass → first order byte; frame occupies TX ~84 cycles"),
    (0, "Settings via CSR map — effect applies at next message boundary"),
    (0, "Tracks lat_min / lat_max / lat_last"),
], size=14, gap=9)
add_rect(slide, Inches(7.0), Inches(1.3), Inches(5.6), Inches(3.1), LIGHT)
add_text(slide, Inches(7.3), Inches(1.45), Inches(5.0), Inches(2.8),
         "2-deep output register, NOT a FIFO\n\n"
         "• 3rd order while 2 pending → dropped, cnt_order_overflow++\n"
         "• Dropping is correct: a queued stale order trades on outdated info\n"
         "• An unbounded queue would make latency history-dependent, "
         "violating NFR-1 (fixed constant latency)",
         size=13, color=DARK)
add_bullets(slide, Inches(0.7), Inches(5.4), Inches(12.0), Inches(1.4), [
    (0, "Orders are never reordered relative to their triggering messages; "
        "market data arriving during TX is still parsed and applied."),
], size=13, gap=8)
footer(slide, 11)

# ---------------------------------------------------------------- slide 12: step 9 stopwatch
slide = content_slide("STEP 9", "The Stopwatch — Proving the Speed")
add_bullets(slide, Inches(0.7), Inches(1.3), Inches(5.2), Inches(2.6), [
    (0, "Free-running 32-bit counter @125 MHz (8 ns period)"),
    (0, "Timestamp IN: last byte of message at parser"),
    (0, "Timestamp OUT: first byte handed to MAC TX"),
    (0, "latency = OUT − IN → 64-bucket BRAM histogram + min/max/last"),
    (0, "Emitting latency inside the order = host capture doubles as a "
        "latency log"),
], size=13, gap=7)
rows = [
    ["Stage", "Cycles"],
    ["Parser ingress → structured message", "1"],
    ["Symbol filter (CAM)", "1"],
    ["Book update (register write)", "1"],
    ["Feature extraction F0–F7", "1"],
    ["Feature normalization", "1"],
    ["ML classifier (parallel MAC)", "2–3"],
    ["ML policy (hysteresis)", "1"],
    ["Risk: all 9 gates incl. ML", "1"],
    ["Order builder → first byte", "1"],
    ["Engine total (tick-to-trade)", "~10–11, target ≤22 (176 ns)"],
]
add_table(slide, Inches(6.2), Inches(1.3), Inches(6.5), 11, 2, rows,
          col_widths=[5.0, 1.5], font_size=11, row_height=0.42, header_size=12)
add_rect(slide, Inches(0.7), Inches(4.2), Inches(5.2), Inches(1.7), LIGHT)
add_text(slide, Inches(1.0), Inches(4.35), Inches(4.6), Inches(1.4),
         "Success = histogram shows EXACTLY ONE occupied bucket over the full "
         "1M-message soak (NFR-1/2). A second bucket is a functional bug, "
         "not a performance result.",
         size=13, color=NAVY, bold=True)
add_text(slide, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.9),
         "Slow path (CSR, stats, debug UART, histogram) never backpressures "
         "the fast path. Single 125 MHz domain; CDC confined to slow path with "
         "two-flop synchronizers.",
         size=12, color=GRAY)
footer(slide, 12)

# ---------------------------------------------------------------- slide 13: success criteria
slide = content_slide(None, "Success Criteria — Project Succeeds If…")
crit = [
    ("1", "Back-to-back messages at Gigabit line rate, no drop/reorder, ≥ 1,000,000 messages"),
    ("2", "Tick-to-trade latency fixed — single histogram bucket, max == min"),
    ("3", "Every risk gate incl. ML 0x09 individually demonstrated blocking/reducing, counter per reason"),
    ("4", "Full RTL output matches Python golden models bit-exactly on randomized soak"),
    ("5", "hls4ml IP score z matches Python fixed-point golden bit-exactly on every regression vector"),
    ("6", "Timing met @125 MHz on XC7A35T (WNS > 0); ≤15% LUTs, ≤10% FFs, ≤8 BRAM, 0 DSP (excl. ML ≤8)"),
    ("7", "One-command reproduce: clone repo → run one command → same simulation results"),
]
box = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.4))
tf = box.text_frame
tf.word_wrap = True
for i, (n, text) in enumerate(crit):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(14)
    r1 = p.add_run()
    r1.text = f"{n}.  "
    r1.font.name = FONT
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT
    r2 = p.add_run()
    r2.text = text
    r2.font.name = FONT
    r2.font.size = Pt(16)
    r2.font.color.rgb = DARK
footer(slide, 13)

# ---------------------------------------------------------------- slide 14: references
slide = content_slide(None, "References & Setup")
add_bullets(slide, Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.3), [
    (-1, "Sources"),
    (0, "fpga_tick_to_trade_master_spec.md v2.0 — single source of truth"),
    (0, "ml_engineer_brief.md — ML owner handoff (Python/hls4ml only, zero Verilog)"),
    (0, "fpga_top_of_book_engine_spec.md — superseded historical document"),
    (-1, "Target"),
    (0, "Artix-7 XC7A35T-2FGG484I, 125 MHz single clock domain"),
    (0, "Vivado 2023.x + Vitis HLS 2023.x"),
    (-1, "Hard rules"),
    (0, "Hand-written RTL is Verilog-2001 only (exception: hls4ml/Vitis IP)"),
    (0, "ML teammate trains in Python (Keras/TF, hls4ml, ap_fixed<8,8> int8), "
        "verifies hls_model.predict() == ml_golden.py bit-exactly before hand-off"),
    (0, "Spec-only, pre-implementation — no rtl/ or sim/ yet"),
], size=14, gap=8)
footer(slide, 14)

# ---------------------------------------------------------------- save
out = r"slides\FPGA_Tick_to_Trade.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
